import os
import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 1. Presentation Initialize (16:9 Widescreen Layout)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 🎨 Professional Corporate White Theme Colors
COLOR_BG = RGBColor(255, 255, 255)         # Pure White Background
COLOR_TEXT_DARK = RGBColor(33, 37, 41)     # Dark Charcoal Text
COLOR_PRIMARY = RGBColor(0, 95, 115)       # Deep Teal Blue (Headings)
COLOR_ACCENT = RGBColor(148, 210, 189)     # Light Teal (Borders)
COLOR_CARD_BG = RGBColor(248, 249, 250)    # Light Gray (Cards)
COLOR_CODE_BG = RGBColor(241, 243, 245)    # Code Block BG
FONT_TITLE = "Arial"
FONT_BODY = "Calibri"

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG

def add_header(slide, title_text, subtitle_text=""):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_TITLE
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = FONT_BODY
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_TEXT_DARK
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.333), Inches(0.02))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT
    shape.line.color.rgb = COLOR_ACCENT

# -------------------------------------------------------------
# SLIDE 1: TITLE SLIDE
# -------------------------------------------------------------
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)

title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.0))
tf = title_box.text_frame
tf.word_wrap = True

p1 = tf.paragraphs[0]
p1.text = "SmartPark Chatbot"
p1.font.name = FONT_TITLE
p1.font.size = Pt(54)
p1.font.bold = True
p1.font.color.rgb = COLOR_PRIMARY

p2 = tf.add_paragraph()
p2.text = "Enterprise RAG-Powered Parking System & Multi-Agent Orchestration"
p2.font.name = FONT_BODY
p2.font.size = Pt(22)
p2.font.color.rgb = COLOR_TEXT_DARK
p2.space_before = Pt(10)

p3 = tf.add_paragraph()
p3.text = "Conversational RAG • LangGraph Workflows • FastAPI MCP • Human-in-the-Loop\nPresenter: Rishabh Gupta"
p3.font.name = FONT_BODY
p3.font.size = Pt(14)
p3.font.color.rgb = COLOR_PRIMARY
p3.space_before = Pt(40)

# -------------------------------------------------------------
# SLIDE 2: 4-STAGE ROADMAP
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Project Delivery Plan", "Phased Architectural Evolution Roadmap")

stages = [
    {"num": "STAGE 1", "title": "Core Engine & CI/CD", "body": "• RAG (Milvus + Azure OpenAI)\n• Intent Orchestration\n• GitHub Actions CI/CD Pipeline"},
    {"num": "STAGE 2", "title": "Persistence & Alerts", "body": "• SQLite Relational Database\n• Static conversation logs\n• Automated Admin Email Alerts"},
    {"num": "STAGE 3", "title": "MCP Integration", "body": "• Python + FastAPI MCP Server\n• Secure File Engine Recording\n• Strict Format Automation"},
    {"num": "STAGE 4", "title": "LangGraph Core", "body": "• Stateful Multi-Agent Graph\n• Human-in-the-Loop Interrupt\n• End-to-End Integration Tests"}
]

card_width = Inches(2.8)
card_height = Inches(4.5)
start_x = Inches(0.5)
y_pos = Inches(1.8)

for i, stage in enumerate(stages):
    x_pos = start_x + (i * Inches(3.1))
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, card_width, card_height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_ACCENT
    tf = card.text_frame
    tf.word_wrap = True
    
    p_num = tf.paragraphs[0]
    p_num.text = stage["num"]
    p_num.font.name = FONT_TITLE; p_num.font.size = Pt(14); p_num.font.bold = True; p_num.font.color.rgb = COLOR_PRIMARY
    
    p_t = tf.add_paragraph()
    p_t.text = stage["title"]
    p_t.font.name = FONT_TITLE; p_t.font.size = Pt(16); p_t.font.bold = True; p_t.font.color.rgb = COLOR_TEXT_DARK; p_t.space_before = Pt(5)
    
    p_b = tf.add_paragraph()
    p_b.text = stage["body"]
    p_b.font.name = FONT_BODY; p_b.font.size = Pt(13); p_b.font.color.rgb = COLOR_TEXT_DARK; p_b.space_before = Pt(15)

# -------------------------------------------------------------
# SLIDE 3: SYSTEM ARCHITECTURE
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "System Architecture & Intent Routing", "LangGraph Stateful Multi-Agent Orchestration Flow")

nodes = [
    {"txt": "User Query Input", "x": 0.5, "y": 2.2, "w": 3.0, "h": 0.8, "is_core": False},
    {"txt": "Input Security Guardrails\n(Checks Injection / Core Vulnerabilities)", "x": 4.0, "y": 2.2, "w": 5.0, "h": 0.8, "is_core": False},
    {"txt": "Node 1: User Interaction & Classifier Router\n(Evaluates Context Intent Prompt via LangGraph)", "x": 0.5, "y": 3.6, "w": 12.333, "h": 0.8, "is_core": True},
    {"txt": "RAG Knowledge Retrieval\n(Milvus + text-embedding-3)", "x": 0.5, "y": 5.0, "w": 3.8, "h": 0.9, "is_core": False},
    {"txt": "Node 2: Admin Approval Node\n(HITL Interrupt / Email Trigger Alerts)", "x": 4.766, "y": 5.0, "w": 3.8, "h": 0.9, "is_core": True},
    {"txt": "Node 3: MCP Storage Recorder\n(FastAPI Server Write Operations)", "x": 9.033, "y": 5.0, "w": 3.8, "h": 0.9, "is_core": True},
    {"txt": "Output Guardrails (Regex PII Masking) ──> Final Structured Clean Response", "x": 0.5, "y": 6.3, "w": 12.333, "h": 0.6, "is_core": False}
]

for n in nodes:
    shape_type = MSO_SHAPE.RECTANGLE if not n["is_core"] else MSO_SHAPE.ROUNDED_RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(n["x"]), Inches(n["y"]), Inches(n["w"]), Inches(n["h"]))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_PRIMARY if n["is_core"] else COLOR_CARD_BG
    box.line.color.rgb = COLOR_ACCENT
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = n["txt"]
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_TITLE if n["is_core"] else FONT_BODY
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255) if n["is_core"] else COLOR_TEXT_DARK

# -------------------------------------------------------------
# SLIDE 4: CORE TECHNOLOGY STACK
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Core Technology Stack", "System Layer Structural Specifications")

rows, cols = 8, 3
left, top, width, height = Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0)
table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(3.5)
table.columns[2].width = Inches(6.333)

headers = ["Layer Name", "Underlying Technology", "Purpose in Architecture"]
table_data = [
    ["Backend Host", "Python 3.11+", "Core programming logic, system processing, pipeline execution"],
    ["RAG Orchestration", "LangChain", "Constructs modular prompt setups, chain flows, connects retrieval units"],
    ["AI Inference", "Azure OpenAI (GPT-4o)", "Context-grounded natural language synthesis and intent classification"],
    ["Embedding Model", "text-embedding-3-small", "Efficient 1536-dimensional semantic vector indexing model"],
    ["Vector Engine", "Milvus (Docker Compose)", "Highly scalable database supporting low-latency vector index lookup"],
    ["CI/CD & Validation", "Pytest + GitHub Actions", "Continuous integration pipelines running quality checks and validation"],
    ["Data Persistence", "SQLite + File Storage", "Structured and unstructured data storage for reservations and logs"]
]

for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.font.name = FONT_TITLE; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = RGBColor(255, 255, 255)

for row_idx, row_data in enumerate(table_data, start=1):
    for col_idx, text in enumerate(row_data):
        cell = table.cell(row_idx, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_CARD_BG if row_idx % 2 == 0 else COLOR_BG
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_BODY; p.font.size = Pt(13); p.font.color.rgb = COLOR_TEXT_DARK
        if col_idx < 2:
            p.font.bold = True

# -------------------------------------------------------------
# SLIDE 5: SETUP GUIDE PART 1 (Cards 01-05)
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Project Setup & Installation Guide (Steps 01-05)", "Sequential System Deployment Cards")

steps_1 = [
    {"id": "01", "name": "Clone Repository", "cmd": "git clone https://github.com/your-org/smartpark-ai.git\ncd smartpark-ai"},
    {"id": "02", "name": "Create Virtual Environment", "cmd": "python3 -m venv .venv\nsource .venv/bin/activate  # Windows: .venv\\Scripts\\activate"},
    {"id": "03", "name": "Install Dependencies", "cmd": "pip install -r requirements.txt"},
    {"id": "04", "name": "Configure Environment Variables", "cmd": "cp .env.example .env\n# Edit with Azure OpenAI & Milvus details"},
    {"id": "05", "name": "Start Milvus Vector Database", "cmd": "cd docker/milvus\ndocker-compose up -d"}
]

card_w = Inches(12.333)
card_h = Inches(0.9)
start_y = Inches(1.8)

for i, step in enumerate(steps_1):
    cur_y = start_y + (i * Inches(1.0))
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), cur_y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_ACCENT
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{step['id']} | {step['name']}  ──>  Code: {step['cmd']}"
    p.font.name = FONT_BODY; p.font.size = Pt(13); p.font.color.rgb = COLOR_TEXT_DARK; p.font.bold = True

# -------------------------------------------------------------
# SLIDE 6: SETUP GUIDE PART 2 (Cards 06-10)
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Project Setup & Installation Guide (Steps 06-10)", "Initialization, Execution and Verification Framework")

steps_2 = [
    {"id": "06", "name": "Initialize DB & Load Docs", "cmd": "python -m app.rag.index_documents\npython -m app.database.init_db"},
    {"id": "07", "name": "Run Chatbot Application", "cmd": "python main.py"},
    {"id": "08", "name": "Run Admin API Gateway", "cmd": "python admin_api.py  # Run in separate terminal window"},
    {"id": "09", "name": "Execute Unit Test Suites", "cmd": "pytest tests/ -v"},
    {"id": "10", "name": "Docker Production Launch", "cmd": "docker-compose -f docker-compose.yml up -d"}
]

for i, step in enumerate(steps_2):
    cur_y = start_y + (i * Inches(1.0))
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), cur_y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = COLOR_ACCENT
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{step['id']} | {step['name']}  ──>  Code: {step['cmd']}"
    p.font.name = FONT_BODY; p.font.size = Pt(13); p.font.color.rgb = COLOR_TEXT_DARK; p.font.bold = True

# -------------------------------------------------------------
# SLIDE 7: CODE ARCHITECTURE
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Code Architecture & Module Mapping", "Production Directory Architecture Map")

box_left = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0))
tf_l = box_left.text_frame
tf_l.word_wrap = True
p = tf_l.paragraphs[0]
p.text = "📂 App Core Infrastructure Routing"
p.font.name = FONT_TITLE; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY

bullets_l = [
    "main.py: Universal entry point for application execution sessions.",
    "chat_orchestrator.py: Routes users between RAG search & Booking Agents.",
    "config.py: Environment variables and configuration systems schema reader.",
    "reservation_agent.py: Collects parameters using Finite State Machine (FSM).",
    "input_filter.py: Runs input guardrail blocks & output masking modules."
]
for b in bullets_l:
    p = tf_l.add_paragraph()
    p.text = "• " + b
    p.font.name = FONT_BODY; p.font.size = Pt(13); p.font.color.rgb = COLOR_TEXT_DARK; p.space_before = Pt(8)

box_right = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(6.0), Inches(5.0))
tf_r = box_right.text_frame
tf_r.word_wrap = True
p_r = tf_r.paragraphs[0]
p_r.text = "📂 RAG Knowledge & Microservices Engine"
p_r.font.name = FONT_TITLE; p_r.font.size = Pt(16); p_r.font.bold = True; p_r.font.color.rgb = COLOR_PRIMARY

bullets_r = [
    "document_loader.py & text_splitter.py: Chunks raw files.",
    "azure_embeddings.py & vector_store.py: Creates embeddings matrix.",
    "milvus_client.py & retriever.py: Executes similarity vector lookups.",
    "rag_chain.py: Merges clean vector contextual data with LLM inference prompts.",
    "admin_api.py / mcp_server.py: FastAPI Model Context Protocol endpoints."
]
for b in bullets_r:
    p = tf_r.add_paragraph()
    p.text = "• " + b
    p.font.name = FONT_BODY; p.font.size = Pt(13); p.font.color.rgb = COLOR_TEXT_DARK; p.space_before = Pt(8)

# -------------------------------------------------------------
# SLIDE 8: RAG ENGINE REFRESH
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "RAG Pipeline & Live Knowledge Base Refreshing", "Zero-Downtime Data Synchronizations Rule Execution")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🔄 Dynamically Refreshing the Vector Knowledge Base:"
p.font.name = FONT_TITLE; p.font.size = Pt(18); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY

p2 = tf.add_paragraph()
p2.text = "When parking rates, regulatory terms or active site locations change, update the environment instantly without stopping the chat servers using:"
p2.font.name = FONT_BODY; p2.font.size = Pt(14); p2.font.color.rgb = COLOR_TEXT_DARK; p2.space_before = Pt(10)

code_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.2), Inches(12.333), Inches(0.8))
code_shape.fill.solid()
code_shape.fill.fore_color.rgb = COLOR_CODE_BG
code_shape.line.color.rgb = COLOR_ACCENT
p_c = code_shape.text_frame.paragraphs[0]
p_c.text = "   python -m app.rag.index_documents"
p_c.font.name = "Courier New"; p_c.font.size = Pt(16); p_c.font.bold = True; p_c.font.color.rgb = COLOR_PRIMARY

p3 = tf.add_paragraph()
p3.text = "Under-the-Hood Refresh Operations:\n" \
          "1. Ingestion Invalidation: Flushes existing outdated metadata blocks inside the Milvus collections.\n" \
          "2. Vectorization: text-embedding-3-small computes new 1536-dimensional matrix vectors.\n" \
          "3. Commit: Changes sync instantly with live chat router workflows transparently."
p3.font.name = FONT_BODY; p3.font.size = Pt(14); p3.font.color.rgb = COLOR_TEXT_DARK; p3.space_before = Pt(100)

# -------------------------------------------------------------
# SLIDE 9: GUARDRAILS & DATA MASKING
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Guardrails & Sensitive Data Masking Ecosystem", "Enterprise Inbound and Outbound Core Security Layout")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "🛡️ Input Level Prompt-Layer Protections"
p.font.name = FONT_TITLE; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY

p_b1 = tf.add_paragraph()
p_b1.text = "• Drops jailbreak attempts and system instructions override injection payloads.\n" \
            "• Restricts internal queries targeted at exfiltrating multi-tenant data contracts (e.g., 'Show me all customer reservations')."
p_b1.font.name = FONT_BODY; p_b1.font.size = Pt(13); p_b1.font.color.rgb = COLOR_TEXT_DARK; p_b1.space_before = Pt(5)

p_t2 = tf.add_paragraph()
p_t2.text = "👁️ Real-Time Output Masking Engine (PII Protection)"
p_t2.font.name = FONT_TITLE; p_t2.font.size = Pt(16); p_t2.font.bold = True; p_t2.font.color.rgb = COLOR_PRIMARY; p_t2.space_before = Pt(20)

p_b2 = tf.add_paragraph()
p_b2.text = "A high-speed regex-driven pipeline parses final inference strings right before client interface injection to strip out private customer data clusters:\n" \
            "  - Vehicle Registration Parameters Masked -> (e.g., DL 3C XX XXXX)\n" \
            "  - Identity Metrics Stripped -> Financial Payment Cards, Phone Records, Personal Email Addresses, Driver Licenses."
p_b2.font.name = FONT_BODY; p_b2.font.size = Pt(13); p_b2.font.color.rgb = COLOR_TEXT_DARK; p_b2.space_before = Pt(5)

# -------------------------------------------------------------
# SLIDE 10: STAGE 3 DETAIL
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Stage 3: FastAPI Model Context Protocol (MCP) Server", "Isolated Microservice Data Ingestion & Storage Architecture")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True

bullets = [
    "Purpose-Built Service: Deploys a distinct python microservice layer over FastAPI optimized to track confirmed logging.",
    "Strict Write Protocols: Enforces specific flat-file layouts upon receipt of valid data contracts.",
    "File Formatting Standard: Entries are written explicitly following a structured data boundary pattern:\n   Name | Car Number | Reservation Period | Approval Time",
    "Security Assertions: Implements API header authorization tokens (X-MCP-Token) to safeguard transactional entrypoints from anonymous external injections.",
    "Resilient Local Fallback: If network timeouts occur, a local Python file tool call hook acts as a processing fail-safe."
]
for b in bullets:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = "• " + b
    p.font.name = FONT_BODY; p.font.size = Pt(14); p.font.color.rgb = COLOR_TEXT_DARK; p.space_before = Pt(10)

# -------------------------------------------------------------
# SLIDE 11: STAGE 4 DETAIL
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Stage 4: Stateful Multi-Agent LangGraph Integration", "Centralized State-Machine Logic Architecture Topology")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True

bullets = [
    "State Engine Architecture: Translates core code elements into manageable nodes with unified schema data contracts passing through edges.",
    "Node 1 - User Interaction & Routing: Executes context parsing, dynamically selecting a RAG lookup or launching reservation tracks.",
    "Node 2 - Administrator Approval (HITL Hook): Implements true Human-in-the-Loop system suspension. Holds state execution in a 'Pending' variable array and triggers admin email notification routines.",
    "Node 3 - Data Recording Gateway: Triggered automatically upon verification to packet structures down to the FastAPI MCP backend.",
    "Stability Assertions: End-to-end conditional routing guarantees predictable graph traversals even across complex parameter inputs."
]
for b in bullets:
    p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
    p.text = "• " + b
    p.font.name = FONT_BODY; p.font.size = Pt(14); p.font.color.rgb = COLOR_TEXT_DARK; p.space_before = Pt(10)

# -------------------------------------------------------------
# SLIDE 12: TESTING & EVALUATION
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "System Verification & High-Throughput Load Testing", "Algorithmic Precision Checks and Edge Failure Load Stress Assessments")

box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(5.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "📉 Algorithmic Evaluation Suites (rag_evaluation.py)"
p.font.name = FONT_TITLE; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = COLOR_PRIMARY

p_b1 = tf.add_paragraph()
p_b1.text = "• Precision Indexing: Evaluates retrieved vector text files to verify alignment and prevent non-relevant background injection.\n" \
            "• Recall Indexing: Confirms no structural compliance or policy details were omitted during the search context assembly step."
p_b1.font.name = FONT_BODY; p_b1.font.size = Pt(13); p_b1.font.color.rgb = COLOR_TEXT_DARK; p_b1.space_before = Pt(5)

p_t2 = tf.add_paragraph()
p_t2.text = "⚡ Stress Verification & High-Volume Operations Analysis"
p_t2.font.name = FONT_TITLE; p_t2.font.size = Pt(16); p_t2.font.bold = True; p_t2.font.color.rgb = COLOR_PRIMARY; p_t2.space_before = Pt(20)

p_b2 = tf.add_paragraph()
p_b2.text = "• Dialogue Mode Integrity: Validates intent classification router processing benchmarks under high parallel query requests.\n" \
            "• HITL Escalation Queue Endurance: Verifies server memory bounds while concurrently holding multiple open booking requests in 'Pending' state.\n" \
            "• FastAPI MCP Concurrency Safety: Stress-tests file locks and SQLite storage layers to guarantee records are written cleanly without data corruption."
p_b2.font.name = FONT_BODY; p_b2.font.size = Pt(13); p_b2.font.color.rgb = COLOR_TEXT_DARK; p_b2.space_before = Pt(5)

# -------------------------------------------------------------
# SLIDE 13: APPENDIX (Stage 3 FastAPI Code)
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Appendix: FastAPI Model Context Protocol (MCP) Server", "Production Endpoint Security and String Writing Structure Code")

code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.4))
tf = code_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "@app.post('/api/mcp/record')\n" \
         "def record_reservation(payload: ReservationSchema, token: str = Depends(verify_mcp_token)):\n" \
         "    try:\n" \
         "        approval_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')\n" \
         "        # Strict Formatting Entry Standard Implemented\n" \
         "        log_entry = f'{payload.name} | {payload.car_number} | {payload.reservation_period} | {approval_time}\\n'\n" \
         "        with open('data/storage/confirmed_reservations.txt', 'a', encoding='utf-8') as f:\n" \
         "            f.write(log_entry)\n" \
         "        return {'status': 'success', 'message': 'Written to storage'}\n" \
         "    except Exception as e:\n" \
         "        raise HTTPException(status_code=500, detail=f'Local Fallback Hook Triggered: {str(e)}')"
p.font.name = "Courier New"; p.font.size = Pt(11); p.font.color.rgb = COLOR_PRIMARY

# -------------------------------------------------------------
# SLIDE 14: APPENDIX (Stage 4 LangGraph Code)
# -------------------------------------------------------------
slide = prs.slides.add_slide(slide_layout)
set_slide_background(slide)
add_header(slide, "Appendix: Multi-Agent State Machine Flow Setup", "LangGraph Structural Composition and Processing Nodes Graph Specification")

code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.4))
tf = code_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "# Unified LangGraph Control Topology Generation\n" \
         "workflow = StateGraph(AgentState)\n" \
         "workflow.add_node('user_interaction', user_interaction_node)\n" \
         "workflow.add_node('admin_approval', admin_approval_node)\n" \
         "workflow.add_node('data_recording', data_recording_node)\n\n" \
         "workflow.set_entry_point('user_interaction')\n" \
         "workflow.add_conditional_edges(\n" \
         "    'user_interaction',\n" \
         "    routing_edge,\n" \
         "    {'admin_approval': 'admin_approval', END: END}\n" \
         ")\n" \
         "workflow.add_edge('admin_approval', 'data_recording')\n" \
         "workflow.add_edge('data_recording', END)\n\n" \
         "smartpark_pipeline = workflow.compile()"
p.font.name = "Courier New"; p.font.size = Pt(11); p.font.color.rgb = COLOR_PRIMARY

# -------------------------------------------------------------
# 💾 TARGET LOCATION CONFIGURATION & SAVE EXECUTION
# -------------------------------------------------------------
output_dir = "/Users/rishabhgupta/Documents/parking-reservation-chatbot/docs/PPTX"
os.makedirs(output_dir, exist_ok=True)  # Automatic folder creation block
final_path = os.path.join(output_dir, "SmartPark_Chatbot_Presentation.pptx")

# Save PPTX File
prs.save(final_path)
print(f"🎉 Success! Your clean white-theme presentation has been saved successfully at:\n{final_path}")